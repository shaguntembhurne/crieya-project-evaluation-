import streamlit as st
import os
import fitz  # PyMuPDF for PDFs
from pptx import Presentation  # For PPTX files
from dotenv import load_dotenv
from reportlab.pdfgen import canvas  # ✅ For actual PDF generation

# ✅ Set page config
st.set_page_config(page_title="CRIEYA Project Evaluation", layout="wide")

# ✅ Load environment variables
load_dotenv()
import google.generativeai as genai

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# ✅ Admin Password (Change this for security)
ADMIN_PASSWORD = "crieya_admin"

# ✅ Google Drive Folder Link (Replace with your actual link)
GOOGLE_DRIVE_FOLDER_LINK = "https://drive.google.com/drive/folders/13noPc-ZIUeUKFGwDjOZGEHPIqi2Hg0ZA?usp=sharing"

# ✅ Function to extract text from PDF
def extract_text_from_pdf(uploaded_file):
    doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    text = "\n".join(page.get_text() for page in doc)
    return text if text.strip() else "No text extracted from PDF."

# ✅ Function to extract text from PPT
def extract_text_from_ppt(uploaded_file):
    prs = Presentation(uploaded_file)
    text = "\n".join(
        [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
    )
    return text if text.strip() else "No extracted text from PPT."

# ✅ Function to generate AI response
def get_gemini_response(evaluation_prompt, specific_prompt, extracted_data):
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content([evaluation_prompt, specific_prompt, extracted_data])
    return response.text

# ✅ Function to save extracted text as an actual PDF
def save_text_as_pdf(text, filename):
    pdf_path = f"{filename}.pdf"
    c = canvas.Canvas(pdf_path)
    c.drawString(100, 800, text[:1000])  # ✅ Only first 1000 chars to avoid overflow
    c.save()
    return pdf_path

# Predefined prompt for scoring and evaluation
evaluation_promt = """You are an AI evaluator for CRIEYA – the Centre for Research, Innovation, and Entrepreneurship for Young Aspirants. CRIEYA supports projects that transform innovative ideas into prototypes and ultimately market-ready products. Your task is to evaluate a submitted project (provided as extracted text from a PDF/PPT) according to the following data-driven criteria and CRIEYA’s strategic guidelines.

1. CRIEYA Mission & Strategic Alignment:
   - Assess whether the project aligns with CRIEYA’s mission to nurture innovation, support prototype development, and drive commercialization.
   - Consider the relevance of the project to key sectors such as Healthcare, Manufacturing & Industry 4.0, Transport & Mobility, Energy & Sustainability, Agriculture & Food Technologies, Media & Entertainment, Defence & Security, Sustainable Design, and Advanced Materials.
   - Does the project reflect the vision of creating breakthrough innovations (e.g., Indian equivalents of major global tech companies)?

2. Funding Context & Justification:
   - Projects may receive up to ₹5,00,000 per project, within an annual innovation fund of ₹3 Crore.
   - Evaluate if the project’s funding request is reasonable and well-justified in relation to its scope, prototype development needs, and the risk de-risking measures proposed.
   - Consider if the proposed budget supports clear milestones from idea to prototype and eventual commercialization.

3. Technology Readiness Level (TRL):
   - Classify the project on a TRL scale (1-9):
     • TRL 1-3: Conceptual stage with early research and proof-of-concept.
     • TRL 4-6: Prototype development and initial real-world testing.
     • TRL 7-9: Market-ready with advanced testing, validation, and readiness for tech transfer.
   - Provide a justification for the TRL classification based on evidence in the submission.

4. Innovation & Novelty:
   - Evaluate the uniqueness of the project. Does it introduce a novel concept, technology, or method?
   - Is the innovation breakthrough in nature, or is it an incremental improvement over existing solutions?
   - Support your evaluation with specifics from the project description.

5. Feasibility & Execution Plan:
   - Analyze the clarity and robustness of the project’s roadmap, including timelines, milestones, and risk management.
   - Are the execution steps practical and realistic to move the project from the idea stage to a prototype and beyond?
   - Identify any potential execution gaps or risks that need addressing.

6. Market Relevance & Impact:
   - Determine whether the project addresses a real-world problem with significant market demand.
   - Is there a clearly defined target market or industry application?
   - Assess the potential economic, societal, or technological impact of the project.

7. Patentability & Intellectual Property (IP) Potential:
   - Evaluate if the project has potential for IP protection and patent filing.
   - Consider whether there is a clear plan for safeguarding intellectual property as the project progresses.

8. Scalability & Commercialization Potential:
   - Assess if the project can scale effectively beyond its current stage.
   - Is there a realistic pathway for technology transfer, startup formation, licensing, or industry partnerships?
   - Consider CRIEYA’s vision of establishing a vibrant innovation ecosystem.

9. Overall Recommendation:
   - Based on the evaluations above, provide a final recommendation:
     • Approve for Funding
     • Needs Further Improvement
     • Reject
   - Include a summary of strengths and weaknesses, as well as key factors that influenced your decision.
   - If improvements are needed, provide clear, actionable recommendations for the applicant.

Your final output should be a structured evaluation report, including:
- A breakdown of scores or ratings (e.g., on a scale of 1-10) for each key category.
- A detailed justification for the TRL classification.
- Specific comments on innovation, feasibility, market relevance, and funding justification.
- A final overall recommendation with a concise rationale.

Use the extracted project document text below for your analysis.
"""
scoring_prompt = """You are an expert evaluator for CRIEYA’s project funding process. Your task is to review the provided project document and generate a structured evaluation report. For each of the categories below, assign a numerical score (on a scale of 1-10, except TRL which is on a scale of 1-9) along with clear justification. Finally, provide an overall rating (1-10) and a final recommendation (Approve, Needs Refinement, or Reject).

1. Innovation & Novelty:
   - Evaluate whether the project introduces a new concept, method, or technology.
   - Rate its uniqueness and potential for breakthrough innovation on a scale of 1-10.
   - Explain the rationale behind the score.

2. Technology Readiness Level (TRL):
   - Classify the project into one of the TRL stages (1 to 9) based on its prototype development, testing, and validation.
   - Rate the readiness on a scale of 1-10, considering that a higher score reflects a more mature, market-ready technology.
   - Provide justification for the classification and score.

3. Market Relevance & Impact:
   - Analyze if the project addresses a significant real-world problem and if there is a strong potential market demand.
   - Rate the market relevance on a scale of 1-10.
   - Support your score with specific observations from the document.

4. Feasibility & Execution Plan:
   - Assess the clarity and robustness of the project’s roadmap, including timelines, milestones, and risk management.
   - Rate the feasibility on a scale of 1-10.
   - Justify your evaluation with examples of strengths or gaps in the execution plan.

5. Funding Justification:
   - Evaluate whether the requested funding (up to ₹5,00,000) is reasonable relative to the project scope.
   - Rate the adequacy of the funding request on a scale of 1-10.
   - Provide specific reasons for the score and comment on budget allocation.

6. Patentability & IP Protection:
   - Determine the project’s potential for intellectual property protection and the likelihood of successful patent filing.
   - Rate this potential on a scale of 1-10.
   - Explain the factors influencing the score.

7. Scalability & Commercialization Potential:
   - Evaluate the project’s ability to scale and transition to a commercially viable product or service.
   - Rate its commercialization potential on a scale of 1-10.
   - Justify your rating with relevant market or strategic insights.

8. Overall Rating & Recommendation:
   - Based on the individual scores, assign an overall rating to the project on a scale of 1-10.
   - Summarize the key strengths and weaknesses.
   - Provide a final recommendation: Approve, Needs Refinement, or Reject, along with a concise explanation.

Use the project document text provided below as your reference for the evaluation.

"""

# Predefined prompt for improvement suggestions
improvement_prompt = """You are an expert evaluator for CRIEYA’s project improvement process. Your task is to analyze the provided project document and generate detailed, actionable feedback that helps the applicant refine their submission. Focus exclusively on identifying areas for improvement and suggesting concrete next steps. Use the project text provided below as your reference.

1. Innovation & Novelty:
   - Evaluate the uniqueness of the project idea.
   - Identify specific gaps or areas where the innovation can be enhanced.
   - Suggest alternative approaches, additional research, or emerging technologies that could make the concept more groundbreaking.

2. Technology Readiness Level (TRL):
   - Assess the current TRL stage of the project and pinpoint what is missing to advance to the next level.
   - Recommend precise technical actions (e.g., further prototype development, enhanced testing protocols, validation studies) that can bolster the project’s maturity.

3. Market Relevance & Impact:
   - Analyze how well the project addresses a real-world problem and meets market demand.
   - Suggest improvements to refine the target market definition or to incorporate additional market research.
   - Recommend partnerships or pilot studies that could better validate the project’s market potential.

4. Feasibility & Execution Plan:
   - Review the project’s implementation roadmap, including timelines and risk management.
   - Identify any shortcomings in planning, such as missing milestones or unaddressed risks.
   - Provide actionable recommendations to streamline the execution plan, optimize timelines, and strengthen risk mitigation strategies.

5. Funding Justification:
   - Examine the alignment between the project scope and the requested funding.
   - Suggest adjustments to the budget or alternative funding strategies (e.g., reallocating costs or seeking complementary grants) to better justify the funding request.

6. Patentability & IP Protection:
   - Evaluate the current approach to intellectual property protection.
   - Recommend measures to improve patentability, such as clearer IP strategy or additional documentation to support patent filings.

7. Scalability & Commercialization:
   - Assess the potential for the project to scale and achieve commercial success.
   - Suggest specific strategies for scaling up, including strategic partnerships, business model improvements, or market entry tactics.
   
8. Overall Recommendations:
   - Summarize the most critical areas for improvement.
   - Provide a prioritized, step-by-step action plan detailing immediate fixes and longer-term enhancements needed for a successful resubmission.

Your final output should be a structured report with clear headings and bullet points for each section, providing the applicant with precise and actionable feedback to improve their project.

"""

# ✅ UI Tabs
tab1, tab2 = st.tabs(["📩 Submit Project", "📊 Evaluate Project"])

# === 📩 SUBMIT PROJECT TAB ===
with tab1:
    st.subheader("Submit Your Project for Evaluation")

    project_title = st.text_input("📌 Project Title")
    project_domain = st.selectbox(
        "🛠 Select Project Domain",
        ["Healthcare", "AI/ML", "Manufacturing", "Energy", "Sustainability", "Other"],
    )
    funding_request = st.number_input("💰 Funding Request (in ₹)", min_value=0, max_value=500000, step=5000)
    trl_level = st.slider("⚙️ Technology Readiness Level (TRL)", min_value=1, max_value=9, value=1)
    upload_file = st.file_uploader("📂 Upload Your Project (PDF/PPTX)", type=["pdf", "pptx"])

    if st.button("Submit Project 🚀"):
        if project_title and project_domain and upload_file:
            file_type = upload_file.type
            extracted_text = extract_text_from_pdf(upload_file) if "pdf" in file_type else extract_text_from_ppt(upload_file)

            if extracted_text:
                # ✅ Save as an actual PDF (Not just a text file)
                pdf_filename = project_title.replace(" ", "_")
                pdf_path = save_text_as_pdf(extracted_text, pdf_filename)

                # ✅ Provide user instructions to manually upload
                
                with open(pdf_path, "rb") as pdf_file:
                    st.download_button(label="📄 Download Extracted PDF", data=pdf_file, file_name=f"{pdf_filename}.pdf", mime="application/pdf")

                st.markdown(f"🔗 **Upload the file manually to your shared Google Drive folder:** [Upload Here]({GOOGLE_DRIVE_FOLDER_LINK})")

                os.remove(pdf_path)  # ✅ Clean up local file after download
        else:
            st.error("⚠️ Please fill in all details and upload a file.")

# === 📊 EVALUATE PROJECT TAB ===
with tab2:
    st.subheader("Evaluate a Project")

    upload_file = st.file_uploader("📂 Upload a Project File for Evaluation", type=["pdf", "pptx"])

    # ✅ Arrange buttons in two columns for better UI
    col1, col2 = st.columns(2)
    with col1:
        submit1 = st.button("🔍 Rate The Project")
    with col2:
        submit2 = st.button("💡 Improvement Suggestions")

    if upload_file is not None:
        file_type = upload_file.type

        if "pdf" in file_type:
            extracted_data = extract_text_from_pdf(upload_file)
        elif "ppt" in file_type or "powerpoint" in file_type:
            extracted_data = extract_text_from_ppt(upload_file)
        else:
            st.error("⚠️ No file uploaded")
            extracted_data = None

        if submit1 and extracted_data:
            response = get_gemini_response(evaluation_promt, scoring_prompt, extracted_data)
            st.subheader("📊 Evaluation Report")
            st.write(response)

        elif submit2 and extracted_data:
            response = get_gemini_response(evaluation_promt, improvement_prompt, extracted_data)
            st.subheader("📈 Improvement Suggestions")
            st.write(response)
    else:
        st.write("📥 Please upload a project file for evaluation.")

